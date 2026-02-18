#!/usr/bin/env python3
"""
Script to convert PowerPoint presentation to HTML
"""

from pptx import Presentation
from datetime import datetime
import re

def escape_html(text):
    """Escape HTML special characters"""
    if not text:
        return ""
    text = str(text)
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;')
    text = text.replace('>', '&gt;')
    text = text.replace('"', '&quot;')
    text = text.replace("'", '&#39;')
    return text

def get_text_from_shape(shape):
    """Extract text from a shape"""
    if hasattr(shape, "text"):
        return shape.text
    if hasattr(shape, "text_frame"):
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            text_parts.append(paragraph.text)
        return "\n".join(text_parts)
    return ""

def get_slide_content(slide):
    """Extract content from a slide"""
    title = ""
    content = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = get_text_from_shape(shape)
            if text.strip():
                # Check if it's a title (usually first shape or title placeholder)
                if shape.shape_type == 1:  # Placeholder type 1 is usually title
                    title = text
                elif not title and len(content) == 0:
                    # First text element is likely the title
                    title = text
                else:
                    content.append(text)
    
    return title, content

def convert_pptx_to_html(pptx_path, output_path):
    """Convert PowerPoint to HTML"""
    prs = Presentation(pptx_path)
    
    html_parts = []
    
    # HTML Header
    html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>InsureMate Project Analysis - Presentation</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            color: #333;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
        }
        
        .container {
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 40px rgba(0,0,0,0.2);
            padding: 40px;
        }
        
        .slide {
            page-break-after: always;
            margin-bottom: 60px;
            padding: 30px;
            background: white;
            border-radius: 8px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            min-height: 600px;
        }
        
        .slide:last-child {
            page-break-after: auto;
        }
        
        .slide-number {
            position: absolute;
            top: 10px;
            right: 20px;
            color: #999;
            font-size: 14px;
        }
        
        h1 {
            color: #667eea;
            text-align: center;
            margin-bottom: 20px;
            font-size: 2.5em;
            border-bottom: 3px solid #667eea;
            padding-bottom: 20px;
        }
        
        h2 {
            color: #764ba2;
            margin-top: 20px;
            margin-bottom: 15px;
            font-size: 2em;
            border-left: 5px solid #667eea;
            padding-left: 15px;
        }
        
        h3 {
            color: #555;
            margin-top: 15px;
            margin-bottom: 10px;
            font-size: 1.5em;
        }
        
        .slide-title {
            color: #764ba2;
            font-size: 2em;
            margin-bottom: 30px;
            padding-bottom: 15px;
            border-bottom: 2px solid #667eea;
        }
        
        .slide-content {
            margin-top: 20px;
        }
        
        .slide-content ul {
            margin-left: 30px;
            margin-top: 15px;
        }
        
        .slide-content li {
            margin: 10px 0;
            font-size: 1.1em;
            line-height: 1.8;
        }
        
        .slide-content p {
            margin: 15px 0;
            font-size: 1.1em;
            line-height: 1.8;
        }
        
        .title-slide {
            text-align: center;
            padding: 80px 40px;
        }
        
        .title-slide h1 {
            font-size: 3em;
            margin-bottom: 30px;
            border: none;
        }
        
        .title-slide .subtitle {
            font-size: 1.5em;
            color: #666;
            margin-top: 30px;
            line-height: 2;
        }
        
        .two-column {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 30px;
            margin-top: 20px;
        }
        
        .column {
            padding: 15px;
        }
        
        .column ul {
            margin-left: 20px;
        }
        
        .footer {
            text-align: center;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 2px solid #eee;
            color: #888;
        }
        
        @media print {
            body {
                background: white;
                padding: 0;
            }
            
            .container {
                box-shadow: none;
                padding: 0;
            }
            
            .slide {
                page-break-after: always;
                margin-bottom: 0;
                box-shadow: none;
            }
        }
    </style>
</head>
<body>
    <div class="container">""")
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        title, content = get_slide_content(slide)
        
        # Determine if it's a title slide (first or last slide)
        is_title_slide = (i == 1 or i == len(prs.slides) or 
                         "Thank You" in title or "Questions" in title)
        
        html_parts.append(f'<div class="slide{" title-slide" if is_title_slide else ""}">')
        html_parts.append(f'<div class="slide-number">Slide {i} of {len(prs.slides)}</div>')
        
        if title:
            if is_title_slide:
                html_parts.append(f'<h1>{escape_html(title)}</h1>')
            else:
                html_parts.append(f'<h2 class="slide-title">{escape_html(title)}</h2>')
        
        if content:
            html_parts.append('<div class="slide-content">')
            for item in content:
                if item.strip():
                    # Check if content has bullet points (lines starting with • or -)
                    lines = item.split('\n')
                    if any(line.strip().startswith('•') or line.strip().startswith('-') 
                           for line in lines if line.strip()):
                        html_parts.append('<ul>')
                        for line in lines:
                            line = line.strip()
                            if line:
                                # Remove bullet markers and add as list item
                                clean_line = re.sub(r'^[•\-\*]\s*', '', line)
                                html_parts.append(f'<li>{escape_html(clean_line)}</li>')
                        html_parts.append('</ul>')
                    else:
                        # Regular paragraph
                        paragraphs = [p for p in item.split('\n') if p.strip()]
                        for para in paragraphs:
                            html_parts.append(f'<p>{escape_html(para)}</p>')
            html_parts.append('</div>')
        
        html_parts.append('</div>')
    
    # HTML Footer
    html_parts.append(f"""        <div class="footer">
            <p><strong>InsureMate Project Analysis - Presentation</strong></p>
            <p>Converted from PowerPoint: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
            <p>Total Slides: {len(prs.slides)}</p>
        </div>
    </div>
</body>
</html>""")
    
    # Write HTML file
    html_content = '\n'.join(html_parts)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return len(prs.slides)

if __name__ == "__main__":
    input_file = "/Users/wekanadmin/InsureMate/InsureMate_Project_Analysis.pptx"
    output_file = "/Users/wekanadmin/InsureMate/InsureMate_Project_Analysis_Presentation.html"
    
    print(f"Converting PowerPoint to HTML...")
    print(f"Input: {input_file}")
    
    try:
        slide_count = convert_pptx_to_html(input_file, output_file)
        print(f"✓ Successfully converted {slide_count} slides")
        print(f"✓ Output: {output_file}")
    except Exception as e:
        print(f"✗ Error: {e}")
        import traceback
        traceback.print_exc()
