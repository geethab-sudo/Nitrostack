#!/usr/bin/env python3
"""
Script to create PowerPoint presentation from project-analysis.html
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from datetime import datetime

def clean_text(text):
    """Remove HTML tags and clean text"""
    # Remove HTML tags
    text = re.sub(r'<[^>]+>', '', text)
    # Remove emojis and special characters that might cause issues
    text = text.replace('🏥', '').replace('📚', '').replace('🏗️', '').replace('📦', '')
    text = text.replace('🛠️', '').replace('🔍', '').replace('💾', '').replace('⚙️', '')
    text = text.replace('🔧', '').replace('🎨', '').replace('📱', '').replace('💼', '')
    text = text.replace('🗄️', '')
    # Clean up whitespace
    text = ' '.join(text.split())
    return text.strip()

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    # Style subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(118, 75, 162)

def add_content_slide(prs, title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = clean_text(title)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(118, 75, 162)
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        if item.strip():
            p = text_frame.add_paragraph()
            p.text = clean_text(item)
            p.level = 0
            p.font.size = Pt(14)
            p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content layout
    title_shape = slide.shapes.title
    left_shape = slide.placeholders[1]
    right_shape = slide.placeholders[2]
    
    title_shape.text = clean_text(title)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(118, 75, 162)
    
    # Left column
    left_frame = left_shape.text_frame
    left_frame.clear()
    for item in left_items:
        if item.strip():
            p = left_frame.add_paragraph()
            p.text = clean_text(item)
            p.level = 0
            p.font.size = Pt(12)
            p.space_after = Pt(4)
    
    # Right column
    right_frame = right_shape.text_frame
    right_frame.clear()
    for item in right_items:
        if item.strip():
            p = right_frame.add_paragraph()
            p.text = clean_text(item)
            p.level = 0
            p.font.size = Pt(12)
            p.space_after = Pt(4)

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    add_title_slide(prs, "InsureMate Project Analysis", 
                    f"Comprehensive Insurance Management System\nTechnical Documentation\n\nGenerated: {datetime.now().strftime('%B %d, %Y')}\nAuthor: Geetha Charu Mathi")
    
    # Technology Stack
    tech_stack = [
        "Runtime & Language: Node.js, TypeScript, ES2022",
        "Framework & Core: NitroStack, @nitrostack/core, @nitrostack/cli",
        "Database: MongoDB (v7.1.0), Connection Pooling",
        "Validation & Schema: Zod (v3.22.4), TypeScript Interfaces",
        "UI Framework: Next.js, React, TypeScript",
        "Utilities: dotenv, In-memory Cache"
    ]
    add_content_slide(prs, "Technology Stack", tech_stack)
    
    # Architecture - Overview
    arch_overview = [
        "5-Layer Architecture",
        "Presentation Layer: Next.js Widgets (React-based UI)",
        "Application Layer: NitroStack MCP Server",
        "Business Logic Layer: Services (Query, Scoring, Recommendation)",
        "Data Access Layer: MongoDBService with Connection Pooling",
        "Data Layer: MongoDB Database with Collections & Indexes"
    ]
    add_content_slide(prs, "System Architecture", arch_overview)
    
    # Architecture - Layers Detail
    arch_layers = [
        "Presentation Layer:",
        "• Next.js Widgets - React-based UI components",
        "• Widgets: insurance-list, insurance-suggestion, insurance-booking, booking-status-list",
        "• Hot module replacement for development",
        "",
        "Application Layer:",
        "• NitroStack MCP Server - Model Context Protocol server",
        "• Tools - Business logic endpoints",
        "• Resources - Data exposure endpoints",
        "• Prompts - Conversation templates"
    ]
    add_content_slide(prs, "Architecture - Presentation & Application Layers", arch_layers)
    
    arch_business = [
        "Business Logic Layer:",
        "• InsuranceTools - Main tool controller",
        "• InsuranceQueryBuilderService - MongoDB query construction",
        "• InsuranceScoringService - Plan matching algorithm",
        "• InsuranceRecommendationService - Recommendation generation",
        "• CacheService - In-memory caching",
        "• MetricsService - Performance monitoring",
        "",
        "Data Access Layer:",
        "• MongoDBService - Database connection management",
        "• Connection Pooling - Efficient connection reuse (min: 2, max: 10)",
        "• Index Management - Automatic index creation",
        "• Retry Logic - Automatic reconnection"
    ]
    add_content_slide(prs, "Architecture - Business Logic & Data Access", arch_business)
    
    # MongoDB Collections
    collections_intro = [
        "Three Main Collections:",
        "1. Insurance Collection (Primary)",
        "2. Users Collection (Secondary)",
        "3. Bookings Collection (Secondary)"
    ]
    add_content_slide(prs, "MongoDB Collections", collections_intro)
    
    insurance_collection = [
        "Insurance Collection:",
        "• Stores all insurance plan information",
        "• Key Fields: name, type, premium, coverage, policyNumber",
        "• Age Range: minAge, maxAge",
        "• Family: familyCoverage, maxFamilyMembers",
        "• Conditions: coversPreExisting, coveredConditions",
        "• Provider & Description",
        "",
        "Indexes:",
        "• Text index on name for search",
        "• Indexes on premium, age ranges, family coverage",
        "• Unique sparse index on policyNumber",
        "• Compound indexes for common queries"
    ]
    add_content_slide(prs, "Insurance Collection", insurance_collection)
    
    users_collection = [
        "Users Collection:",
        "• Stores user account information",
        "• Key Fields: email (unique), name, phone, age, salary",
        "• Additional: designation, address, familyMembers",
        "• Timestamps: createdAt, updatedAt",
        "",
        "Indexes:",
        "• Unique index on email (primary key)",
        "• Text index on name for search",
        "• Indexes on age, salary, designation",
        "• Sparse unique index on phone",
        "• Index on createdAt for sorting"
    ]
    add_content_slide(prs, "Users Collection", users_collection)
    
    bookings_collection = [
        "Bookings Collection:",
        "• Stores insurance booking records",
        "• References: userId, insurancePlanId, policyNumber",
        "• Status: status (pending/confirmed/cancelled/expired/active)",
        "• Financial: premium, coverageAmount",
        "• Dates: startDate, endDate",
        "• Payment: paymentStatus, paymentMethod, transactionId",
        "• Additional: notes, createdAt, updatedAt",
        "",
        "Indexes:",
        "• Index on userId (most common query)",
        "• Indexes on insurancePlanId, status, paymentStatus",
        "• Compound indexes for user bookings by status",
        "• Date range indexes for queries"
    ]
    add_content_slide(prs, "Bookings Collection", bookings_collection)
    
    # Tools & Usage
    tools_intro = [
        "Four Main Tools:",
        "1. list_insurance - Query Tool",
        "2. suggest_insurance_plan - AI Tool",
        "3. book_insurance_with_user - Transaction Tool",
        "4. list_booking_status_by_email - Query Tool"
    ]
    add_content_slide(prs, "Tools & Usage", tools_intro)
    
    list_insurance = [
        "list_insurance Tool:",
        "Purpose: List and filter insurance plans",
        "Widget: insurance-list",
        "",
        "Input Parameters:",
        "• salary (Required) - Annual salary for affordability",
        "• familyMembers (Optional) - Number of family members",
        "• age (Optional, 0-120) - User age for age range matching",
        "• limit (Optional, default: 100) - Max records",
        "• skip (Optional, default: 0) - Pagination",
        "• filter (Optional) - Additional MongoDB filter",
        "",
        "Filtering Logic:",
        "• Premium ≤ 10% of annual salary",
        "• Family coverage matching",
        "• Age range validation",
        "• Custom MongoDB operators supported"
    ]
    add_content_slide(prs, "Tool: list_insurance", list_insurance)
    
    suggest_insurance = [
        "suggest_insurance_plan Tool:",
        "Purpose: AI-powered insurance plan suggestions with match scoring",
        "Widget: insurance-suggestion",
        "",
        "Input Parameters:",
        "• age (Required, 0-120) - User age",
        "• salary (Required) - Annual salary",
        "• familyMembers (Required) - Number of family members",
        "• deficiencies (Optional) - Health deficiencies/pre-existing conditions",
        "• insuranceType (Optional) - Filter by type",
        "",
        "Scoring Algorithm (0-100 points):",
        "• Age Match: 30 points",
        "• Premium Affordability: 25 points",
        "• Family Coverage: 20 points",
        "• Pre-existing Conditions: 25 points",
        "• Specific Conditions: 15 points",
        "• Type Match: 10 points",
        "",
        "Returns: Top 5 plans sorted by match score"
    ]
    add_content_slide(prs, "Tool: suggest_insurance_plan", suggest_insurance)
    
    book_insurance = [
        "book_insurance_with_user Tool:",
        "Purpose: Create booking record and manage user account",
        "Widget: insurance-booking",
        "",
        "Input Parameters:",
        "• policyNumber (Required) - Insurance policy number",
        "• email (Required) - User email (validated)",
        "• name (Required) - User full name",
        "• phoneNumber (Required) - User phone number",
        "• paymentMethod (Optional) - Payment method",
        "• startDate (Optional) - ISO 8601 format",
        "• years (Optional, 1-50) - Coverage duration",
        "• transactionId (Optional) - Payment transaction ID",
        "• notes (Optional) - Additional notes",
        "",
        "Business Logic:",
        "• Validates insurance plan exists",
        "• Creates/updates user account",
        "• Creates booking with status='pending'",
        "• Sets payment status based on transactionId",
        "• Calculates endDate from startDate + years"
    ]
    add_content_slide(prs, "Tool: book_insurance_with_user", book_insurance)
    
    list_booking = [
        "list_booking_status_by_email Tool:",
        "Purpose: Retrieve all bookings for a user by email",
        "Widget: booking-status-list",
        "",
        "Input Parameters:",
        "• email (Required) - User email address",
        "",
        "Business Logic:",
        "• Validates email format",
        "• Finds user by email (returns error if not found)",
        "• Retrieves all bookings for the user",
        "• Sorts by createdAt descending (most recent first)",
        "",
        "Response:",
        "• success - Boolean",
        "• user - User document",
        "• bookings - Array of all booking documents",
        "• count - Total number of bookings"
    ]
    add_content_slide(prs, "Tool: list_booking_status_by_email", list_booking)
    
    # MongoDB Queries
    queries = [
        "MongoDB Query Optimization:",
        "• InsuranceQueryBuilderService constructs optimized queries",
        "• All queries leverage database indexes",
        "• Filters applied at database level",
        "• Pagination with skip/limit",
        "• Projection for necessary fields only",
        "• Input sanitization prevents injection attacks",
        "",
        "Query Types:",
        "1. List Insurance Query - Filters by salary, age, family",
        "2. Suggestion Query - Optimized for scoring algorithm",
        "3. Booking Queries - User lookup, plan lookup, status filtering"
    ]
    add_content_slide(prs, "MongoDB Queries", queries)
    
    # Cache Handling
    cache_intro = [
        "Cache Implementation:",
        "• Type: In-memory Map-based cache",
        "• Pattern: Singleton service",
        "• Storage: JavaScript Map<string, CacheEntry>",
        "• TTL: Configurable per entry, default 5 minutes",
        "",
        "Cache Operations:",
        "• get<T>(key) - Retrieve cached value",
        "• set<T>(key, value, ttl?) - Store value",
        "• getOrCompute<T>(key, computeFn, ttl?) - Cache or compute",
        "• delete(key) - Remove entry",
        "• clear() - Remove all entries"
    ]
    add_content_slide(prs, "Cache Handling", cache_intro)
    
    cache_details = [
        "Cache TTL by Operation:",
        "• list_insurance: 5 minutes",
        "• suggest_insurance_plan: 5 minutes",
        "• search_insurance_names: 2 minutes",
        "",
        "Cache Key Generation:",
        "• Deterministic from input parameters",
        "• Parameters sorted alphabetically",
        "• Values JSON stringified",
        "• Format: 'prefix:param1:value1|param2:value2|...'",
        "",
        "Cache Cleanup:",
        "• Automatic cleanup every 60 seconds",
        "• Lazy expiration on get() operations",
        "• Removes expired entries"
    ]
    add_content_slide(prs, "Cache Details", cache_details)
    
    cache_benefits = [
        "Cache Benefits:",
        "• Performance: Reduces database queries by 60-80%",
        "• Latency: Cache hits return in <1ms vs 10-50ms for DB queries",
        "• Database Load: Significantly reduces load on MongoDB",
        "• Cost: Lower database operation costs",
        "",
        "Cache Flow:",
        "1. Request received with parameters",
        "2. Generate deterministic cache key",
        "3. Check cache for existing entry",
        "4. If cache hit and not expired, return cached data",
        "5. If cache miss, execute database query",
        "6. Store result in cache with TTL",
        "7. Return result to caller"
    ]
    add_content_slide(prs, "Cache Benefits & Flow", cache_benefits)
    
    # How It Works
    system_init = [
        "System Initialization:",
        "1. Application Start - NitroStack CLI starts MCP server",
        "2. Module Loading - AppModule loads InsuranceModule",
        "3. MongoDB Initialization - Connection established",
        "4. Connection Pooling - Pool created (min: 2, max: 10)",
        "5. Index Creation - Automatic index creation for collections",
        "6. Service Initialization - CacheService, MetricsService as singletons",
        "7. Health Check Start - Monitoring every 30 seconds",
        "8. Ready State - Server ready to accept requests"
    ]
    add_content_slide(prs, "System Initialization", system_init)
    
    request_flow = [
        "Request Processing Flow:",
        "1. Request Received - MCP server receives tool invocation",
        "2. Input Validation - Zod schema validates parameters",
        "3. Input Sanitization - InputSanitizer sanitizes inputs",
        "4. Cache Check - CacheService checks for cached result",
        "5. Query Building - InsuranceQueryBuilderService builds query",
        "6. Database Query - MongoDBService executes using indexes",
        "7. Data Processing - Results processed (scoring, formatting)",
        "8. Cache Storage - Result stored in cache with TTL",
        "9. Metrics Recording - MetricsService records operation",
        "10. Response Return - Formatted response returned"
    ]
    add_content_slide(prs, "Request Processing Flow", request_flow)
    
    error_handling = [
        "Error Handling:",
        "• Custom Error Classes with codes and context",
        "• Error Types:",
        "  - DatabaseConnectionError - Connection failures",
        "  - DatabaseQueryError - Query execution failures",
        "  - InvalidInputError - Input validation failures",
        "  - ConfigurationError - Configuration issues",
        "",
        "Retry Logic:",
        "• Automatic reconnection with exponential backoff",
        "• Maximum 5 retry attempts",
        "• All errors tracked in MetricsService"
    ]
    add_content_slide(prs, "Error Handling", error_handling)
    
    performance = [
        "Performance Optimizations:",
        "• Database Indexes: 20+ indexes across collections",
        "• Connection Pooling: Reuses connections (2-10)",
        "• Caching: In-memory cache reduces database load",
        "• Query Optimization: Filters at database level",
        "• Pagination: Efficient data retrieval with skip/limit",
        "• Lazy Loading: Collections initialized when needed",
        "",
        "Monitoring & Metrics:",
        "• MetricsService tracks: operation count, avg/min/max times",
        "• Error count and success rate percentage",
        "• Health Checks: MongoDB connection status, ping every 30s"
    ]
    add_content_slide(prs, "Performance & Monitoring", performance)
    
    security = [
        "Security Features:",
        "• Input Sanitization:",
        "  - Regex injection prevention",
        "  - String sanitization",
        "  - Number validation",
        "  - Filter object sanitization",
        "",
        "• Type Validation:",
        "  - Zod schemas validate all inputs",
        "",
        "• MongoDB Injection Prevention:",
        "  - Filter objects sanitized before query",
        "",
        "• Email Validation:",
        "  - Regex validation for email format"
    ]
    add_content_slide(prs, "Security Features", security)
    
    # Services & Components
    core_services = [
        "Core Services:",
        "• MongoDBService - Database connection, pooling, indexes (Singleton)",
        "• CacheService - In-memory caching with TTL (Singleton)",
        "• MetricsService - Operation tracking, monitoring (Singleton)",
        "• InsuranceQueryBuilderService - MongoDB query construction",
        "• InsuranceScoringService - Match score calculation",
        "• InsuranceRecommendationService - Recommendation generation"
    ]
    add_content_slide(prs, "Core Services", core_services)
    
    configuration = [
        "Configuration:",
        "• InsuranceConfig: Environment variable management with Zod",
        "",
        "Environment Variables:",
        "• MONGODB_URI - MongoDB connection string",
        "• MONGODB_DATABASE_NAME - Database name (default: 'Insurance')",
        "• MONGODB_COLLECTION_NAME - Collection name (default: 'Insurance')",
        "• MONGODB_MAX_POOL_SIZE - Max connections (default: 10)",
        "• MONGODB_MIN_POOL_SIZE - Min connections (default: 2)",
        "• MONGODB_CONNECT_TIMEOUT_MS - Connection timeout (default: 30000)",
        "• MONGODB_SERVER_SELECTION_TIMEOUT_MS - Server selection timeout (default: 5000)"
    ]
    add_content_slide(prs, "Configuration", configuration)
    
    utilities = [
        "Utilities:",
        "• InputSanitizer - Sanitizes strings, numbers, arrays, filter objects",
        "• ObjectIdUtil - Converts MongoDB ObjectIds to strings for JSON"
    ]
    add_content_slide(prs, "Utilities", utilities)
    
    # Widgets
    widgets = [
        "Widgets (UI Components):",
        "Next.js-based React widgets for displaying tool results:",
        "",
        "• insurance-list - Displays list of insurance plans",
        "• insurance-suggestion - Shows suggested plans with scores",
        "• insurance-booking - Booking confirmation interface",
        "• booking-status-list - User booking history with status cards",
        "• calculator-result - Calculator tool result display",
        "• insurance-search-dropdown - Search dropdown component",
        "",
        "Widget Development:",
        "• Hot module replacement enabled for development"
    ]
    add_content_slide(prs, "Widgets (UI Components)", widgets)
    
    # Summary Slide
    summary = [
        "Key Highlights:",
        "• Modern Tech Stack: Node.js, TypeScript, MongoDB, Next.js",
        "• 5-Layer Architecture with clear separation of concerns",
        "• 3 MongoDB Collections with optimized indexes",
        "• 4 Main Tools for insurance management",
        "• Intelligent Caching System (60-80% query reduction)",
        "• Comprehensive Error Handling & Retry Logic",
        "• Performance Optimizations (20+ indexes, connection pooling)",
        "• Security Features (Input sanitization, injection prevention)",
        "• Monitoring & Metrics for operational insights"
    ]
    add_content_slide(prs, "Summary", summary)
    
    # Thank You Slide
    add_title_slide(prs, "Thank You", "Questions & Discussion")
    
    return prs

if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    prs = create_presentation()
    output_file = "/Users/wekanadmin/InsureMate/InsureMate_Project_Analysis.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
